"""
报告生成模块 - PDF/Excel/HTML/PPT报告生成 + 模板选择

功能:
1. 分析报告生成
2. 多格式导出 (PDF, Excel, HTML, PPT)
3. 模板选择 (professional, academic, executive)
4. 历史报告管理
"""

import os
import json
from datetime import datetime
import pandas as pd


# ============================================================================
# 报告模板定义
# ============================================================================

TEMPLATES = {
    'professional': {
        'name': '专业版',
        'description': '标准专业风格，适合日常分析报告',
        'primary_color': '#667eea',
        'secondary_color': '#1E3A5F',
        'accent_color': '#764ba2',
        'bg_color': '#f5f5f5',
        'text_color': '#333333',
        'title_font_size': 18,
        'heading_font_size': 14,
        'body_font_size': 10,
        'table_header_bg': '#667eea',
        'disclaimer_bg': '#fff3cd',
        'disclaimer_border': '#ffc107',
    },
    'academic': {
        'name': '学术版',
        'description': '严谨学术风格，适合论文和学术交流',
        'primary_color': '#2c3e50',
        'secondary_color': '#34495e',
        'accent_color': '#7f8c8d',
        'bg_color': '#ecf0f1',
        'text_color': '#2c3e50',
        'title_font_size': 16,
        'heading_font_size': 13,
        'body_font_size': 10,
        'table_header_bg': '#2c3e50',
        'disclaimer_bg': '#e8e8e8',
        'disclaimer_border': '#7f8c8d',
    },
    'executive': {
        'name': '高管简版',
        'description': '简洁高管风格，突出关键指标和结论',
        'primary_color': '#e74c3c',
        'secondary_color': '#c0392b',
        'accent_color': '#f39c12',
        'bg_color': '#fdf2e9',
        'text_color': '#2c3e50',
        'title_font_size': 22,
        'heading_font_size': 16,
        'body_font_size': 12,
        'table_header_bg': '#e74c3c',
        'disclaimer_bg': '#fef9e7',
        'disclaimer_border': '#f39c12',
    },
}


# ============================================================================
# P0修复: 免责声明
# ============================================================================

DISCLAIMER = """
⚠️ 重要声明 (Disclaimer)

本报告仅供学术研究和教育目的，不构成任何投资建议。

1. 模型局限性：所有计量经济学模型都是对现实的简化，实际市场行为可能偏离模型预测。
2. 历史不代表未来：基于历史数据的统计特征不保证在未来延续。
3. 风险提示：投资有风险，决策需谨慎。请在专业人士指导下做出投资决策。
4. 免责：报告作者不对因使用本报告内容而导致的任何损失承担责任。

本报告基于CNLocalGovSpread框架生成，版本: 3.0.0
"""


# ============================================================================
# 报告生成器
# ============================================================================

class ReportGenerator:
    """报告生成器类"""

    def __init__(self, output_dir='reports'):
        """
        初始化报告生成器

        参数:
            output_dir: 报告输出目录
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        self.report_history_file = os.path.join(output_dir, 'history.json')
        self.template = TEMPLATES['professional']
        self._load_history()

    def _load_history(self):
        """加载报告历史"""
        if os.path.exists(self.report_history_file):
            with open(self.report_history_file, 'r') as f:
                self.history = json.load(f)
        else:
            self.history = []

    def _save_history(self):
        """保存报告历史"""
        with open(self.report_history_file, 'w') as f:
            json.dump(self.history, f, indent=2, ensure_ascii=False)

    def generate_report(self, clean_data, returns, kalman, vol_modeler, evt,
                            title="地方债利差分析报告", format="PDF", sections=None,
                            template='professional'):
        """
        生成分析报告

        参数:
            clean_data: 清洗后的数据
            returns: 收益率序列
            kalman: 卡尔曼滤波器
            vol_modeler: 波动率建模器
            evt: EVT分析器
            title: 报告标题
            format: 报告格式 (PDF, Excel, HTML, PPT)
            sections: 包含的章节列表
            template: 报告模板 (professional, academic, executive)

        返回:
            str: 报告文件路径
        """
        if sections is None:
            sections = ['数据概览', '信号分析', '波动率分析', '风险分析', '交易建议']

        if template not in TEMPLATES:
            raise ValueError(f"不支持的模板: {template}, 可选: {list(TEMPLATES.keys())}")
        self.template = TEMPLATES[template]

        # 生成报告内容
        report_data = self._prepare_report_data(
            clean_data, returns, kalman, vol_modeler, evt, sections
        )

        # 根据格式生成报告
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        if format.upper() == 'PDF':
            path = self._generate_pdf(report_data, title, timestamp)
        elif format.upper() == 'EXCEL':
            path = self._generate_excel(report_data, title, timestamp)
        elif format.upper() == 'HTML':
            path = self._generate_html(report_data, title, timestamp)
        elif format.upper() == 'PPT':
            path = self._generate_ppt(report_data, title, timestamp)
        else:
            raise ValueError(f"不支持的格式: {format}, 可选: PDF, Excel, HTML, PPT")

        # 记录历史
        self._add_to_history(title, format, path, sections)
        return path

    def _prepare_report_data(self, clean_data, returns, kalman, vol_modeler, evt, sections):
        """准备报告数据"""
        data = {}
        spread = clean_data['spread']

        # 数据概览
        if '数据概览' in sections:
            data['overview'] = {
                '数据范围': f"{clean_data.index[0].strftime('%Y-%m-%d')} 至 {clean_data.index[-1].strftime('%Y-%m-%d')}",
                '交易日数': len(clean_data),
                '当前利差': f"{spread.iloc[-1]:.4f}",
                '历史均值': f"{spread.mean():.4f}",
                '历史标准差': f"{spread.std():.4f}",
                '最小值': f"{spread.min():.4f}",
                '最大值': f"{spread.max():.4f}",
                '偏度': f"{spread.skew():.4f}",
                '峰度': f"{spread.kurt():.4f}"
            }

        # 信号分析
        if '信号分析' in sections and kalman is not None:
            smoothed = kalman.smoothed_state
            deviation = kalman.get_signal_deviation()
            current_dev = deviation.iloc[-1]

            signal = '中性'
            if current_dev > 1.5:
                signal = '做空（利差高估）'
            elif current_dev < -1.5:
                signal = '做多（利差低估）'

            data['signal'] = {
                '当前利差': f"{spread.iloc[-1]:.4f}",
                '趋势水平': f"{smoothed.iloc[-1]:.4f}",
                '偏离度': f"{current_dev:.4f}σ",
                '交易信号': signal
            }

        # 波动率分析
        if '波动率分析' in sections and vol_modeler is not None:
            winner = vol_modeler.run_tournament()
            winner_vol = vol_modeler.get_conditional_volatility(winner)

            data['volatility'] = {
                '获胜模型': winner,
                'AIC': f"{vol_modeler.ic_scores[winner]['AIC']:.2f}",
                'BIC': f"{vol_modeler.ic_scores[winner]['BIC']:.2f}",
                '当前波动率': f"{winner_vol.iloc[-1]:.6f}",
                '平均波动率': f"{winner_vol.mean():.6f}"
            }

        # 风险分析
        if '风险分析' in sections and evt is not None:
            var = evt.calculate_var() if evt.var is None else evt.var
            es = evt.calculate_es() if evt.es is None else evt.es
            tail_idx = evt.get_tail_index()

            data['risk'] = {
                '99% VaR': f"{var:.6f}" if var else "N/A",
                '99% ES': f"{es:.6f}" if es else "N/A",
                '尾部指数': f"{tail_idx:.4f}" if tail_idx else "N/A",
                'GPD形状参数': f"{evt.gpd_params['shape']:.4f}" if evt.gpd_params else "N/A",
                'GPD尺度参数': f"{evt.gpd_params['scale']:.4f}" if evt.gpd_params else "N/A"
            }

        # 交易建议
        if '交易建议' in sections:
            data['recommendation'] = self._generate_recommendation(
                clean_data, kalman, evt, vol_modeler
            )

        return data

    def _generate_recommendation(self, clean_data, kalman, evt, vol_modeler):
        """生成交易建议"""
        spread = clean_data['spread']
        current = spread.iloc[-1]
        mean = spread.mean()

        # 综合判断
        signals = []

        # 偏离度信号
        if kalman is not None:
            deviation = kalman.get_signal_deviation().iloc[-1]
            if deviation > 1.5:
                signals.append(('做空', deviation, '利差高估'))
            elif deviation < -1.5:
                signals.append(('做多', deviation, '利差低估'))
            else:
                signals.append(('中性', deviation, '利差正常'))

        # 趋势信号
        recent_mean = spread.iloc[-20:].mean()
        if recent_mean > mean:
            signals.append(('观望', 0, '上升趋势'))
        elif recent_mean < mean:
            signals.append(('观望', 0, '下降趋势'))

        # 综合建议
        recommendation = {
            '当前利差': f"{current:.4f}",
            '历史均值': f"{mean:.4f}",
            '建议': signals[0][0] if signals else '中性',
            '理由': signals[0][2] if signals else '数据正常',
            '止损建议': f"{current + (evt.var if evt and evt.var else spread.std()):.4f}"
        }

        return recommendation

    def _generate_pdf(self, data, title, timestamp):
        """生成PDF报告"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib import colors
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            tmpl = self.template

            # 注册中文字体（跨平台支持）
            chinese_font = 'Helvetica'

            font_paths = [
                '/System/Library/Fonts/PingFang.ttc',
                '/System/Library/Fonts/STHeiti Light.ttc',
                '/Library/Fonts/Arial Unicode.ttf',
                'C:/Windows/Fonts/simhei.ttf',
                'C:/Windows/Fonts/msyh.ttc',
                '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc',
                '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
            ]

            for font_path in font_paths:
                try:
                    if os.path.exists(font_path):
                        pdfmetrics.registerFont(TTFont('SimHei', font_path))
                        chinese_font = 'SimHei'
                        break
                except Exception:
                    continue

            filename = os.path.join(self.output_dir, f"report_{timestamp}.pdf")
            doc = SimpleDocTemplate(filename, pagesize=A4)

            # 模板驱动的颜色解析
            primary_hex = tmpl['primary_color'].lstrip('#')
            primary_rgb = colors.Color(
                int(primary_hex[0:2], 16)/255,
                int(primary_hex[2:4], 16)/255,
                int(primary_hex[4:6], 16)/255
            )
            header_hex = tmpl['table_header_bg'].lstrip('#')
            header_rgb = colors.Color(
                int(header_hex[0:2], 16)/255,
                int(header_hex[2:4], 16)/255,
                int(header_hex[4:6], 16)/255
            )

            styles = getSampleStyleSheet()
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontName=chinese_font,
                fontSize=tmpl['title_font_size'],
                spaceAfter=30,
                textColor=primary_rgb
            )
            heading_style = ParagraphStyle(
                'CustomHeading',
                parent=styles['Heading2'],
                fontName=chinese_font,
                fontSize=tmpl['heading_font_size'],
                spaceAfter=12,
                textColor=primary_rgb
            )
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontName=chinese_font,
                fontSize=tmpl['body_font_size'],
            )

            elements = []

            elements.append(Paragraph(title, title_style))
            elements.append(Paragraph(
                f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}  |  模板: {tmpl['name']}",
                body_style
            ))
            elements.append(Spacer(1, 1*cm))

            for section, content in data.items():
                elements.append(Paragraph(section, heading_style))

                if isinstance(content, dict):
                    table_data = [[k, v] for k, v in content.items()]
                    table = Table(table_data, colWidths=[6*cm, 8*cm])
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (0, -1), colors.lightgrey),
                        ('BACKGROUND', (0, 0), (-1, 0), header_rgb),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                        ('FONTNAME', (0, 0), (-1, -1), chinese_font),
                        ('FONTSIZE', (0, 0), (-1, -1), tmpl['body_font_size']),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('LEFTPADDING', (0, 0), (-1, -1), 6),
                    ]))
                    elements.append(table)
                    elements.append(Spacer(1, 0.5*cm))

            doc.build(elements)
            return filename

        except ImportError:
            return self._generate_text(data, title, timestamp, 'txt')

    def _generate_excel(self, data, title, timestamp):
        """生成Excel报告"""
        filename = os.path.join(self.output_dir, f"report_{timestamp}.xlsx")

        with pd.ExcelWriter(filename, engine='openpyxl') as writer:
            # 摘要页
            summary_data = []
            for section, content in data.items():
                if isinstance(content, dict):
                    for key, value in content.items():
                        summary_data.append({
                            '章节': section,
                            '指标': key,
                            '值': value
                        })

            if summary_data:
                summary_df = pd.DataFrame(summary_data)
                summary_df.to_excel(writer, sheet_name='报告摘要', index=False)

            # 各章节单独页
            for section, content in data.items():
                if isinstance(content, dict):
                    df = pd.DataFrame([content]).T
                    df.columns = ['值']
                    df.index.name = '指标'
                    df.to_excel(writer, sheet_name=section[:31])  # Excel sheet名最长31字符

        return filename

    def _generate_html(self, data, title, timestamp):
        """生成HTML报告"""
        tmpl = self.template
        filename = os.path.join(self.output_dir, f"report_{timestamp}.html")

        html = f"""
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
            background: {tmpl['bg_color']};
            color: {tmpl['text_color']};
        }}
        .container {{
            background: white;
            padding: 30px;
            border-radius: 10px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        h1 {{
            color: {tmpl['secondary_color']};
            border-bottom: 3px solid {tmpl['primary_color']};
            padding-bottom: 10px;
            font-size: {tmpl['title_font_size']}px;
        }}
        h2 {{
            color: {tmpl['primary_color']};
            margin-top: 30px;
            font-size: {tmpl['heading_font_size']}px;
        }}
        .meta {{
            color: #666;
            font-size: 0.9rem;
            margin-bottom: 20px;
        }}
        .template-badge {{
            display: inline-block;
            background: {tmpl['primary_color']};
            color: white;
            padding: 2px 8px;
            border-radius: 4px;
            font-size: 0.75rem;
            margin-left: 8px;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 15px 0;
            font-size: {tmpl['body_font_size']}px;
        }}
        th, td {{
            padding: 12px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }}
        th {{
            background: {tmpl['table_header_bg']};
            color: white;
        }}
        tr:hover {{
            background: #f5f5f5;
        }}
        .footer {{
            margin-top: 40px;
            padding-top: 20px;
            border-top: 1px solid #ddd;
            text-align: center;
            color: #666;
            font-size: 0.8rem;
        }}
        .disclaimer {{
            margin-top: 30px;
            padding: 20px;
            background: {tmpl['disclaimer_bg']};
            border-left: 4px solid {tmpl['disclaimer_border']};
            font-size: 0.85rem;
            color: #856404;
        }}
    </style>
</head>
<body>
    <div class="container">
        <h1>{title} <span class="template-badge">{tmpl['name']}</span></h1>
        <p class="meta">生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        <div class="disclaimer">
            <strong>⚠️ 重要声明</strong><br>
            本报告仅供学术研究和教育目的，不构成任何投资建议。所有模型都是对现实的简化，历史表现不代表未来收益。投资有风险，决策需谨慎。
        </div>
"""

        for section, content in data.items():
            html += f"        <h2>{section}</h2>\n"
            html += "        <table>\n"
            html += "            <tr><th>指标</th><th>值</th></tr>\n"
            if isinstance(content, dict):
                for key, value in content.items():
                    html += f"            <tr><td>{key}</td><td>{value}</td></tr>\n"
            html += "        </table>\n"

        html += f"""
        <div class="footer">
            <p>CNLocalGovSpread v3.0.0 | Author: Quinn Liu</p>
            <p><a href="https://github.com/quinnmacro/CNLocalGovSpread">GitHub</a> | <a href="https://www.linkedin.com/in/liulu-math">LinkedIn</a></p>
        </div>
    </div>
</body>
</html>
"""

        with open(filename, 'w', encoding='utf-8') as f:
            f.write(html)

        return filename

    def _generate_ppt(self, data, title, timestamp):
        """生成PPT报告"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            tmpl = self.template

            def hex_to_rgb(hex_str):
                h = hex_str.lstrip('#')
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            primary = hex_to_rgb(tmpl['primary_color'])
            secondary = hex_to_rgb(tmpl['secondary_color'])
            accent = hex_to_rgb(tmpl['accent_color'])
            text_color = hex_to_rgb(tmpl['text_color'])

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # ---- 标题幻灯片 ----
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(tmpl['title_font_size'])
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = secondary
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            subtitle = slide.placeholders[1]
            subtitle.text = f"CNLocalGovSpread v3.0.0 | {tmpl['name']}模板\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            for para in subtitle.text_frame.paragraphs:
                para.font.size = Pt(14)
                para.font.color.rgb = primary

            # ---- 关键指标摘要幻灯片 (executive模板专属) ----
            if tmpl == TEMPLATES['executive'] and data:
                summary_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(summary_layout)
                slide.shapes.title.text = "关键指标摘要"
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(tmpl['heading_font_size'])
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary

                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                key_items = []
                for section_name, content in data.items():
                    if isinstance(content, dict):
                        for k, v in content.items():
                            if k in ('当前利差', '建议', '交易信号', '99% VaR',
                                     '获胜模型', '偏离度', '尾部指数'):
                                key_items.append(f"{k}: {v}")

                for i, item in enumerate(key_items[:8]):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = item
                    para.font.size = Pt(tmpl['body_font_size'] + 2)
                    para.font.color.rgb = text_color
                    para.font.bold = True if i < 4 else False
                    if para.font.bold:
                        para.font.color.rgb = accent

            # ---- 各章节幻灯片 ----
            for section, content in data.items():
                content_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = section
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(tmpl['heading_font_size'])
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary

                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                if isinstance(content, dict):
                    for i, (key, value) in enumerate(content.items()):
                        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        para.text = f"{key}: {value}"
                        para.font.size = Pt(tmpl['body_font_size'])
                        para.font.color.rgb = text_color
                        para.space_after = Pt(4)

            # ---- 免责声明幻灯片 ----
            disclaimer_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(disclaimer_layout)
            slide.shapes.title.text = "免责声明"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(tmpl['heading_font_size'])
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            disclaimer_lines = [
                "本报告仅供学术研究和教育目的，不构成任何投资建议。",
                "所有计量经济学模型都是对现实的简化，实际市场行为可能偏离模型预测。",
                "基于历史数据的统计特征不保证在未来延续。",
                "投资有风险，决策需谨慎。请在专业人士指导下做出投资决策。",
                f"CNLocalGovSpread v3.0.0 | Author: Quinn Liu"
            ]
            for i, line in enumerate(disclaimer_lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = line
                para.font.size = Pt(tmpl['body_font_size'])
                para.font.color.rgb = text_color

            filename = os.path.join(self.output_dir, f"report_{timestamp}.pptx")
            prs.save(filename)
            return filename

        except ImportError:
            return self._generate_text(data, title, timestamp, 'txt')

    def _generate_text(self, data, title, timestamp, ext='txt'):
        """生成文本报告"""
        filename = os.path.join(self.output_dir, f"report_{timestamp}.{ext}")

        lines = [
            "=" * 60,
            title,
            "=" * 60,
            f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            ""
        ]

        for section, content in data.items():
            lines.append(f"\n{section}")
            lines.append("-" * 40)
            if isinstance(content, dict):
                for key, value in content.items():
                    lines.append(f"  {key}: {value}")

        lines.append("\n" + "=" * 60)

        with open(filename, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))

        return filename

    def _add_to_history(self, title, format, path, sections):
        """添加到历史记录"""
        record = {
            'title': title,
            'format': format,
            'path': path,
            'sections': sections,
            'timestamp': datetime.now().isoformat(),
            'filename': os.path.basename(path)
        }
        self.history.insert(0, record)
        self._save_history()

    def get_history(self):
        """获取报告历史"""
        return self.history

    def delete_report(self, filename):
        """删除报告"""
        path = os.path.join(self.output_dir, filename)
        if os.path.exists(path):
            os.remove(path)
            self.history = [r for r in self.history if r.get('filename') != filename]
            self._save_history()
            return True
        return False


# ============================================================================
# 便捷函数
# ============================================================================

def generate_report(clean_data, returns, kalman, vol_modeler, evt,
                        title="地方债利差分析报告", format="PDF", sections=None,
                        template='professional'):
    """
    生成分析报告的便捷函数

    参数:
        clean_data: 清洗后的数据
        returns: 收益率序列
        kalman: 卡尔曼滤波器
        vol_modeler: 波动率建模器
        evt: EVT分析器
        title: 报告标题
        format: 报告格式 (PDF, Excel, HTML, PPT)
        sections: 包含的章节
        template: 报告模板 (professional, academic, executive)

    返回:
        str: 报告文件路径
    """
    generator = ReportGenerator()
    return generator.generate_report(
        clean_data, returns, kalman, vol_modeler, evt,
        title=title, format=format, sections=sections, template=template
    )


def get_report_history():
    """获取报告历史的便捷函数"""
    generator = ReportGenerator()
    history = generator.get_history()

    if not history:
        return None

    return pd.DataFrame(history)


def generate_quick_report(clean_data, returns, kalman, vol_modeler, evt):
    """生成快速报告（包含所有章节）"""
    return generate_report(
        clean_data, returns, kalman, vol_modeler, evt,
        title="地方债利差快速分析报告",
        format="HTML",
        sections=['数据概览', '信号分析', '波动率分析', '风险分析', '交易建议']
    )
