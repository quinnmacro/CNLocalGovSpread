"""
report_gen.py 测试 - 报告生成模块

覆盖 ReportGenerator 类和3个便捷函数:
- 初始化/历史加载/保存
- generate_report (PDF/Excel/HTML/Text 格式)
- _prepare_report_data (各章节数据准备)
- _generate_recommendation (交易建议生成)
- _add_to_history / get_history / delete_report
- generate_report / get_report_history / generate_quick_report 便捷函数
- DISCLAIMER 常量
"""

import pytest
import os
import json
import numpy as np
import pandas as pd
from unittest.mock import MagicMock, patch
from datetime import datetime

from src.report_gen import (
    ReportGenerator,
    DISCLAIMER,
    TEMPLATES,
    generate_report,
    get_report_history,
    generate_quick_report,
)


# ============================================================================
# Fixtures
# ============================================================================

@pytest.fixture
def tmp_output_dir(tmp_path):
    """临时报告输出目录"""
    return str(tmp_path / 'reports')


@pytest.fixture
def clean_data():
    """模拟清洗后的数据"""
    dates = pd.date_range('2020-01-01', periods=100, freq='D')
    np.random.seed(42)
    spread = 80 + np.cumsum(np.random.randn(100) * 0.3)
    return pd.DataFrame({'spread': spread}, index=dates)


@pytest.fixture
def returns():
    """模拟收益率序列"""
    dates = pd.date_range('2020-01-02', periods=99, freq='D')
    np.random.seed(42)
    return pd.Series(np.random.randn(99) * 0.02, index=dates, name='returns')


@pytest.fixture
def mock_kalman():
    """模拟卡尔曼滤波器"""
    kalman = MagicMock()
    dates = pd.date_range('2020-01-01', periods=100, freq='D')
    kalman.smoothed_state = pd.Series(80 + np.random.randn(100) * 0.5, index=dates)
    kalman.get_signal_deviation.return_value = pd.Series(np.random.randn(100), index=dates)
    return kalman


@pytest.fixture
def mock_vol_modeler():
    """模拟波动率建模器"""
    vol = MagicMock()
    vol.run_tournament.return_value = 'GARCH'
    vol.ic_scores = {'GARCH': {'AIC': -500.0, 'BIC': -480.0}}
    dates = pd.date_range('2020-01-01', periods=100, freq='D')
    vol.get_conditional_volatility.return_value = pd.Series(
        np.abs(np.random.randn(100)) * 0.02 + 0.01, index=dates
    )
    return vol


@pytest.fixture
def mock_evt():
    """模拟EVT分析器"""
    evt = MagicMock()
    evt.var = -0.05
    evt.es = -0.08
    evt.calculate_var.return_value = -0.05
    evt.calculate_es.return_value = -0.08
    evt.get_tail_index.return_value = 0.3
    evt.gpd_params = {'shape': 0.15, 'scale': 0.02}
    return evt


@pytest.fixture
def generator(tmp_output_dir):
    """创建报告生成器实例"""
    return ReportGenerator(output_dir=tmp_output_dir)


# ============================================================================
# DISCLAIMER
# ============================================================================

class TestDisclaimer:

    def test_disclaimer_exists(self):
        assert DISCLAIMER is not None
        assert len(DISCLAIMER) > 0

    def test_disclaimer_contains_version(self):
        assert '3.0.0' in DISCLAIMER

    def test_disclaimer_contains_warning(self):
        assert '不构成任何投资建议' in DISCLAIMER

    def test_disclaimer_contains_risk_notice(self):
        assert '投资有风险' in DISCLAIMER


# ============================================================================
# ReportGenerator.__init__
# ============================================================================

class TestReportGeneratorInit:

    def test_creates_output_dir(self, tmp_output_dir):
        gen = ReportGenerator(output_dir=tmp_output_dir)
        assert os.path.isdir(tmp_output_dir)

    def test_default_output_dir(self):
        gen = ReportGenerator()
        assert gen.output_dir == 'reports'

    def test_history_file_path(self, tmp_output_dir):
        gen = ReportGenerator(output_dir=tmp_output_dir)
        expected = os.path.join(tmp_output_dir, 'history.json')
        assert gen.report_history_file == expected

    def test_empty_history_on_new_dir(self, tmp_output_dir):
        gen = ReportGenerator(output_dir=tmp_output_dir)
        assert gen.history == []


# ============================================================================
# _load_history / _save_history
# ============================================================================

class TestHistoryManagement:

    def test_load_history_from_file(self, tmp_output_dir):
        os.makedirs(tmp_output_dir, exist_ok=True)
        history = [{'title': 'test', 'format': 'PDF', 'path': '/tmp/test.pdf'}]
        with open(os.path.join(tmp_output_dir, 'history.json'), 'w') as f:
            json.dump(history, f)

        gen = ReportGenerator(output_dir=tmp_output_dir)
        assert len(gen.history) == 1
        assert gen.history[0]['title'] == 'test'

    def test_save_history_creates_file(self, tmp_output_dir):
        gen = ReportGenerator(output_dir=tmp_output_dir)
        gen.history = [{'title': 'saved', 'format': 'HTML'}]
        gen._save_history()

        with open(os.path.join(tmp_output_dir, 'history.json'), 'r') as f:
            loaded = json.load(f)
        assert loaded[0]['title'] == 'saved'

    def test_get_history_returns_list(self, generator):
        result = generator.get_history()
        assert isinstance(result, list)


# ============================================================================
# _prepare_report_data
# ============================================================================

class TestPrepareReportData:

    def test_overview_section(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览']
        )
        assert 'overview' in data
        assert '当前利差' in data['overview']
        assert '历史均值' in data['overview']

    def test_signal_section(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['信号分析']
        )
        assert 'signal' in data
        assert '交易信号' in data['signal']

    def test_signal_without_kalman(self, generator, clean_data, returns, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, None, mock_vol_modeler, mock_evt,
            sections=['信号分析']
        )
        assert 'signal' not in data

    def test_volatility_section(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['波动率分析']
        )
        assert 'volatility' in data
        assert '获胜模型' in data['volatility']

    def test_risk_section(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['风险分析']
        )
        assert 'risk' in data
        assert '99% VaR' in data['risk']

    def test_recommendation_section(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['交易建议']
        )
        assert 'recommendation' in data
        assert '建议' in data['recommendation']

    def test_all_sections(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览', '信号分析', '波动率分析', '风险分析', '交易建议']
        )
        assert len(data) >= 4


# ============================================================================
# _generate_recommendation
# ============================================================================

class TestGenerateRecommendation:

    def test_with_kalman_high_deviation(self, generator, clean_data, mock_kalman):
        # Set deviation > 1.5 (高估 → 做空)
        mock_kalman.get_signal_deviation.return_value = pd.Series([2.0])
        rec = generator._generate_recommendation(clean_data, mock_kalman, None, None)
        assert rec['建议'] == '做空'

    def test_with_kalman_low_deviation(self, generator, clean_data, mock_kalman):
        # Set deviation < -1.5 (低估 → 做多)
        mock_kalman.get_signal_deviation.return_value = pd.Series([-2.0])
        rec = generator._generate_recommendation(clean_data, mock_kalman, None, None)
        assert rec['建议'] == '做多'

    def test_with_kalman_normal_deviation(self, generator, clean_data, mock_kalman):
        mock_kalman.get_signal_deviation.return_value = pd.Series([0.5])
        rec = generator._generate_recommendation(clean_data, mock_kalman, None, None)
        assert rec['建议'] == '中性'

    def test_without_kalman(self, generator, clean_data):
        rec = generator._generate_recommendation(clean_data, None, None, None)
        assert '建议' in rec
        assert '当前利差' in rec

    def test_stop_loss_with_evt_var(self, generator, clean_data, mock_evt):
        mock_evt.var = -0.05
        rec = generator._generate_recommendation(clean_data, None, mock_evt, None)
        assert '止损建议' in rec

    def test_stop_loss_without_evt(self, generator, clean_data):
        rec = generator._generate_recommendation(clean_data, None, None, None)
        assert '止损建议' in rec


# ============================================================================
# generate_report (format dispatch)
# ============================================================================

class TestGenerateReportFormats:

    def test_html_format(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        assert path.endswith('.html')
        assert os.path.exists(path)

    def test_html_contains_title(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            title="测试报告", format="HTML"
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '测试报告' in content

    def test_html_contains_disclaimer(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '重要声明' in content

    def test_html_contains_version(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert 'v3.0.0' in content

    def test_excel_format(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="Excel"
        )
        assert path.endswith('.xlsx')
        assert os.path.exists(path)

    def test_text_format_fallback(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        # TXT is not a supported format - should raise ValueError
        with pytest.raises(ValueError, match='不支持的格式'):
            generator.generate_report(
                clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                format="TXT"
            )

    def test_unsupported_format_raises(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with pytest.raises(ValueError, match='不支持的格式'):
            generator.generate_report(
                clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                format="DOCX"
            )

    def test_custom_sections(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML", sections=['数据概览']
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert 'overview' in content

    def test_default_sections(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert 'overview' in content

    def test_pdf_format_or_text_fallback(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PDF"
        )
        # Either PDF or text fallback depending on reportlab availability
        assert os.path.exists(path)


# ============================================================================
# _generate_html internal
# ============================================================================

class TestGenerateHtmlInternal:

    def test_html_structure(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览']
        )
        path = generator._generate_html(data, "Test Report", "20200101_120000")
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '<!DOCTYPE html>' in content
        assert '</html>' in content

    def test_html_tables(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览']
        )
        path = generator._generate_html(data, "Test", "20200101")
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '<table>' in content


# ============================================================================
# _generate_text internal
# ============================================================================

class TestGenerateTextInternal:

    def test_text_format(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览']
        )
        path = generator._generate_text(data, "Test Report", "20200101_120000", 'txt')
        assert path.endswith('.txt')
        assert os.path.exists(path)

    def test_text_content(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        data = generator._prepare_report_data(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            sections=['数据概览']
        )
        path = generator._generate_text(data, "Test Report", "20200101", 'txt')
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert 'Test Report' in content
        assert 'overview' in content


# ============================================================================
# _add_to_history / get_history / delete_report
# ============================================================================

class TestHistoryOps:

    def test_add_to_history(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        history = generator.get_history()
        assert len(history) == 1
        assert history[0]['format'] == 'HTML'

    def test_multiple_reports_in_history(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        history = generator.get_history()
        assert len(history) == 2

    def test_delete_report(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        filename = os.path.basename(path)
        result = generator.delete_report(filename)
        assert result is True
        assert not os.path.exists(path)
        assert len(generator.get_history()) == 0

    def test_delete_nonexistent_report(self, generator):
        result = generator.delete_report('nonexistent.html')
        assert result is False

    def test_history_record_fields(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            title="测试标题", format="HTML"
        )
        record = generator.get_history()[0]
        assert 'title' in record
        assert 'format' in record
        assert 'path' in record
        assert 'timestamp' in record
        assert 'filename' in record
        assert record['title'] == '测试标题'


# ============================================================================
# Convenience functions
# ============================================================================

class TestConvenienceFunctions:

    def test_generate_report_function(self, tmp_path, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with patch.object(ReportGenerator, '__init__', return_value=None):
            with patch.object(ReportGenerator, 'generate_report', return_value='/tmp/test.html') as mock_gen:
                result = generate_report(clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt)
                assert mock_gen.called

    def test_get_report_history_function(self, tmp_output_dir):
        gen = ReportGenerator(output_dir=tmp_output_dir)
        # Empty history → None
        with patch('src.report_gen.ReportGenerator') as mock_class:
            mock_instance = MagicMock()
            mock_instance.get_history.return_value = []
            mock_class.return_value = mock_instance
            result = get_report_history()
            assert result is None

    def test_generate_quick_report_function(self, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with patch.object(ReportGenerator, '__init__', return_value=None):
            with patch.object(ReportGenerator, 'generate_report', return_value='/tmp/quick.html') as mock_gen:
                result = generate_quick_report(clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt)
                assert mock_gen.called


# ============================================================================
# TEMPLATES
# ============================================================================

class TestTemplates:

    def test_templates_dict_exists(self):
        assert isinstance(TEMPLATES, dict)
        assert len(TEMPLATES) >= 3

    def test_professional_template(self):
        assert 'professional' in TEMPLATES
        t = TEMPLATES['professional']
        assert t['name'] == '专业版'
        assert t['primary_color'].startswith('#')
        assert 'title_font_size' in t

    def test_academic_template(self):
        assert 'academic' in TEMPLATES
        t = TEMPLATES['academic']
        assert t['name'] == '学术版'
        assert 'table_header_bg' in t

    def test_executive_template(self):
        assert 'executive' in TEMPLATES
        t = TEMPLATES['executive']
        assert t['name'] == '高管简版'
        assert t['title_font_size'] >= TEMPLATES['professional']['title_font_size']

    def test_template_required_keys(self):
        required_keys = ['name', 'description', 'primary_color', 'secondary_color',
                         'accent_color', 'bg_color', 'text_color', 'title_font_size',
                         'heading_font_size', 'body_font_size', 'table_header_bg',
                         'disclaimer_bg', 'disclaimer_border']
        for key, tmpl in TEMPLATES.items():
            for rk in required_keys:
                assert rk in tmpl, f"Template '{key}' missing key '{rk}'"

    def test_template_colors_are_hex(self):
        color_keys = ['primary_color', 'secondary_color', 'accent_color',
                      'bg_color', 'text_color', 'table_header_bg',
                      'disclaimer_bg', 'disclaimer_border']
        for key, tmpl in TEMPLATES.items():
            for ck in color_keys:
                assert tmpl[ck].startswith('#'), f"Template '{key}' color '{ck}' not hex"


# ============================================================================
# Template selection in generate_report
# ============================================================================

class TestTemplateSelection:

    def test_default_template_is_professional(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML"
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '专业版' in content

    def test_academic_template_in_html(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML", template='academic'
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '学术版' in content
        assert '#2c3e50' in content

    def test_executive_template_in_html(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML", template='executive'
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '高管简版' in content
        assert '#e74c3c' in content

    def test_invalid_template_raises(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with pytest.raises(ValueError, match='不支持的模板'):
            generator.generate_report(
                clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                format="HTML", template='nonexistent'
            )

    def test_template_stored_on_generator(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML", template='academic'
        )
        assert generator.template == TEMPLATES['academic']

    def test_template_applies_font_sizes(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="HTML", template='executive'
        )
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        assert '22px' in content  # executive title_font_size

    def test_template_in_convenience_function(self, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with patch.object(ReportGenerator, '__init__', return_value=None):
            with patch.object(ReportGenerator, 'generate_report', return_value='/tmp/test.html') as mock_gen:
                result = generate_report(
                    clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                    template='academic'
                )
                call_kwargs = mock_gen.call_args[1]
                assert call_kwargs['template'] == 'academic'


# ============================================================================
# PPT generation
# ============================================================================

class TestPPTGeneration:

    def test_ppt_format_generates_file(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT"
        )
        assert path.endswith('.pptx')
        assert os.path.exists(path)

    def test_ppt_has_multiple_slides(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", sections=['数据概览', '风险分析']
        )
        prs = Presentation(path)
        # Title slide + 2 section slides + disclaimer slide = 4
        assert len(prs.slides) >= 4

    def test_ppt_title_slide_content(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", title="测试PPT报告"
        )
        prs = Presentation(path)
        title_text = prs.slides[0].shapes.title.text
        assert '测试PPT报告' == title_text

    def test_ppt_subtitle_has_version(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT"
        )
        prs = Presentation(path)
        subtitle_text = prs.slides[0].placeholders[1].text
        assert 'v3.0.0' in subtitle_text

    def test_ppt_disclaimer_slide(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", sections=['数据概览']
        )
        prs = Presentation(path)
        # Last slide should be disclaimer
        last_slide = prs.slides[-1]
        title_text = last_slide.shapes.title.text
        assert '免责声明' == title_text

    def test_ppt_with_academic_template(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", template='academic'
        )
        prs = Presentation(path)
        subtitle_text = prs.slides[0].placeholders[1].text
        assert '学术版' in subtitle_text

    def test_ppt_executive_has_summary_slide(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", template='executive',
            sections=['数据概览', '信号分析', '风险分析']
        )
        prs = Presentation(path)
        # Title + executive summary + 3 sections + disclaimer = 6
        assert len(prs.slides) >= 6

    def test_ppt_section_slide_content(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        path = generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT", sections=['数据概览']
        )
        prs = Presentation(path)
        # Slide 1: title, Slide 2: 数据概览, Slide 3: disclaimer
        section_slide = prs.slides[1]
        title_text = section_slide.shapes.title.text
        assert '数据概览' == title_text or 'overview' == title_text

    def test_ppt_pptx_import_fallback(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        # When python-pptx is not available, should fall back to text
        with patch.dict('sys.modules', {'pptx': None}):
            data = generator._prepare_report_data(
                clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                sections=['数据概览']
            )
            generator.template = TEMPLATES['professional']
            path = generator._generate_ppt(data, "Test", "20200101")
            assert path.endswith('.txt')

    def test_ppt_history_records_format(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        generator.generate_report(
            clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
            format="PPT"
        )
        record = generator.get_history()[0]
        assert record['format'] == 'PPT'
        assert record['path'].endswith('.pptx')

    def test_ppt_format_in_error_message(self, generator, clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt):
        with pytest.raises(ValueError, match='PDF, Excel, HTML, PPT'):
            generator.generate_report(
                clean_data, returns, mock_kalman, mock_vol_modeler, mock_evt,
                format="DOCX"
            )